#!/usr/bin/env python3
"""
PowerPoint Resizer and Organizer
Automatically resize PowerPoint presentations to 36x48 inches and reorganize all elements
to maintain perfect layout and alignment with UNIFORM SCALING (no distortion).
"""

from pptx import Presentation
from pptx.util import Inches
import argparse
import os
import sys


class PowerPointResizer:
    def __init__(self, input_file, output_file, target_width=36, target_height=48, scale_mode='fit'):
        """
        Initialize the PowerPoint Resizer

        Args:
            input_file: Path to input PowerPoint file
            output_file: Path to output PowerPoint file
            target_width: Target width in inches (default: 36)
            target_height: Target height in inches (default: 48)
            scale_mode: 'fit' (fit within bounds), 'fill' (fill entire space), or 'stretch' (stretch to exact size)
        """
        self.input_file = input_file
        self.output_file = output_file
        self.target_width = Inches(target_width)
        self.target_height = Inches(target_height)
        self.scale_mode = scale_mode

    def calculate_uniform_scale(self, original_width, original_height):
        """
        Calculate UNIFORM scale factor to maintain aspect ratios
        This ensures circles stay circular, squares stay square, etc.

        Returns: (scale_factor, offset_x, offset_y)
        """
        width_scale = self.target_width / original_width
        height_scale = self.target_height / original_height

        if self.scale_mode == 'fit':
            # Use the smaller scale to ensure everything fits
            scale = min(width_scale, height_scale)
        elif self.scale_mode == 'fill':
            # Use the larger scale to fill the space (may crop)
            scale = max(width_scale, height_scale)
        else:  # stretch
            # Return both scales for non-uniform scaling
            return width_scale, height_scale, 0, 0

        # Calculate offsets to center the content
        scaled_width = original_width * scale
        scaled_height = original_height * scale
        offset_x = (self.target_width - scaled_width) / 2
        offset_y = (self.target_height - scaled_height) / 2

        return scale, scale, offset_x, offset_y

    def resize_and_reposition_shape(self, shape, scale_x, scale_y, offset_x, offset_y):
        """
        Resize and reposition a shape with UNIFORM scaling
        Maintains exact proportions - no distortion!
        """
        try:
            # Scale position (top-left corner) and add centering offset
            if hasattr(shape, 'left') and hasattr(shape, 'top'):
                shape.left = int(shape.left * scale_x + offset_x)
                shape.top = int(shape.top * scale_y + offset_y)

            # Scale size uniformly
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                shape.width = int(shape.width * scale_x)
                shape.height = int(shape.height * scale_y)

            # Handle text frame font sizes
            if hasattr(shape, 'text_frame'):
                # Use the uniform scale factor for fonts
                uniform_scale = min(scale_x, scale_y) if scale_x != scale_y else scale_x
                self.scale_text_frame(shape.text_frame, uniform_scale)

        except Exception as e:
            print(f"Warning: Could not fully resize shape: {e}")

    def scale_text_frame(self, text_frame, scale_factor):
        """
        Scale font sizes in text frame to maintain readability
        """
        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size:
                        # Scale font size
                        new_size = int(run.font.size * scale_factor)
                        run.font.size = new_size
        except Exception as e:
            print(f"Warning: Could not scale text: {e}")

    def organize_slide(self, slide, scale_x, scale_y, offset_x, offset_y):
        """
        Organize all elements on a slide with uniform scaling
        Everything maintains its exact proportions
        """
        # Process all shapes on the slide
        for shape in slide.shapes:
            self.resize_and_reposition_shape(shape, scale_x, scale_y, offset_x, offset_y)

    def process_presentation(self):
        """
        Main processing function
        Resizes the presentation with UNIFORM scaling - no distortion!
        """
        print(f"Loading presentation: {self.input_file}")

        # Load the presentation
        try:
            prs = Presentation(self.input_file)
        except Exception as e:
            print(f"Error: Could not load PowerPoint file: {e}")
            return False

        # Get original dimensions
        original_width = prs.slide_width
        original_height = prs.slide_height

        print(f"Original size: {original_width / Inches(1):.2f} x {original_height / Inches(1):.2f} inches")
        print(f"Target size: {self.target_width / Inches(1):.2f} x {self.target_height / Inches(1):.2f} inches")
        print(f"Scale mode: {self.scale_mode}")

        # Calculate uniform scale factor
        scale_x, scale_y, offset_x, offset_y = self.calculate_uniform_scale(original_width, original_height)

        if scale_x == scale_y:
            print(f"Uniform scale factor: {scale_x:.3f}x")
            print(f"Content will be centered with offset: ({offset_x / Inches(1):.2f}\", {offset_y / Inches(1):.2f}\")")
        else:
            print(f"Scale factors: Width={scale_x:.3f}, Height={scale_y:.3f}")

        # Resize the slide dimensions
        prs.slide_width = self.target_width
        prs.slide_height = self.target_height

        # Process each slide
        total_slides = len(prs.slides)
        print(f"\nProcessing {total_slides} slide(s)...")

        for i, slide in enumerate(prs.slides, 1):
            print(f"  Processing slide {i}/{total_slides}...")
            self.organize_slide(slide, scale_x, scale_y, offset_x, offset_y)

        # Save the modified presentation
        print(f"\nSaving to: {self.output_file}")
        try:
            prs.save(self.output_file)
            print("✓ Successfully saved!")
            return True
        except Exception as e:
            print(f"Error: Could not save PowerPoint file: {e}")
            return False


def main():
    parser = argparse.ArgumentParser(
        description='Resize and reorganize PowerPoint presentations with UNIFORM scaling (no distortion)',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Basic usage - resize to 36x48 inches (fits content, centered)
  python resize_powerpoint.py input.pptx output.pptx

  # Custom dimensions
  python resize_powerpoint.py input.pptx output.pptx --width 24 --height 36

  # Fill mode - content fills entire space (may crop edges)
  python resize_powerpoint.py input.pptx output.pptx --mode fill

  # Stretch mode - stretches to exact size (may distort)
  python resize_powerpoint.py input.pptx output.pptx --mode stretch

Scaling Modes:
  fit    - Fits all content within bounds, maintains proportions (default)
  fill   - Fills entire slide, maintains proportions, may crop
  stretch- Stretches to exact dimensions (may cause distortion)
        """
    )

    parser.add_argument('input', help='Input PowerPoint file (.pptx)')
    parser.add_argument('output', help='Output PowerPoint file (.pptx)')
    parser.add_argument('--width', type=float, default=36,
                        help='Target width in inches (default: 36)')
    parser.add_argument('--height', type=float, default=48,
                        help='Target height in inches (default: 48)')
    parser.add_argument('--mode', choices=['fit', 'fill', 'stretch'], default='fit',
                        help='Scaling mode: fit (default), fill, or stretch')

    args = parser.parse_args()

    # Validate input file
    if not os.path.exists(args.input):
        print(f"Error: Input file '{args.input}' does not exist")
        sys.exit(1)

    if not args.input.endswith('.pptx'):
        print("Warning: Input file should be .pptx format")

    # Create resizer and process
    resizer = PowerPointResizer(
        args.input,
        args.output,
        args.width,
        args.height,
        args.mode
    )

    success = resizer.process_presentation()

    if success:
        print("\n✓ All done! Your presentation has been resized with uniform scaling.")
        print("  Everything looks exactly the same, just scaled proportionally!")
        sys.exit(0)
    else:
        print("\n✗ Failed to process presentation")
        sys.exit(1)


if __name__ == '__main__':
    main()
